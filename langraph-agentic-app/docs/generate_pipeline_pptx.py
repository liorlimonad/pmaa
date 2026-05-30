"""
Generate pipeline-diagram.pptx  —  academic block-flow diagram of the
Process Mining & Agentic Pipeline (pmaa workspace).

Run from repo root:
    source .venv/bin/activate
    python langraph-agentic-app/docs/generate_pipeline_pptx.py
"""

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ── Slide dimensions: 16:9 widescreen ───────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.50)

# ── Colour palette ───────────────────────────────────────────────────────────
C_BG        = RGBColor(0xFA, 0xFA, 0xF8)
C_TITLETXT  = RGBColor(0x11, 0x11, 0x11)
C_SUBTXT    = RGBColor(0x44, 0x44, 0x44)
C_IDXTXT    = RGBColor(0x88, 0x88, 0x88)
C_LINE      = RGBColor(0x44, 0x44, 0x44)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# node theme colours (fill, border)
THEME = {
    "source":   (RGBColor(0xDC, 0xE8, 0xF5), RGBColor(0x4A, 0x7D, 0xB5)),
    "process":  (RGBColor(0xEE, 0xF3, 0xE8), RGBColor(0x5A, 0x8A, 0x4A)),
    "output":   (RGBColor(0xFD, 0xF5, 0xE0), RGBColor(0xC0, 0x90, 0x20)),
    "mining":   (RGBColor(0xF0, 0xEB, 0xF8), RGBColor(0x7A, 0x5A, 0xAD)),
    "llm":      (RGBColor(0xFD, 0xEE, 0xF0), RGBColor(0xC0, 0x50, 0x60)),
    "agent":    (RGBColor(0xE8, 0xF4, 0xF4), RGBColor(0x2A, 0x8A, 0x8A)),
    "subagent": (RGBColor(0xC8, 0xE8, 0xE8), RGBColor(0x2A, 0x8A, 0x8A)),
    "eval":     (RGBColor(0xF4, 0xF4, 0xF4), RGBColor(0x77, 0x77, 0x77)),
}


# ─────────────────────────────────────────────────────────────────────────────
# Primitive helpers
# ─────────────────────────────────────────────────────────────────────────────

def _rgb_hex(rgb):
    return "{:02X}{:02X}{:02X}".format(rgb[0], rgb[1], rgb[2])


def add_rect(slide, x, y, w, h, theme_key, border_pt=1.4, round_pct=20000):
    fill_rgb, border_rgb = THEME[theme_key]
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(border_pt)
    sp = shape.element
    spPr = sp.find(qn("p:spPr"))
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "roundRect")
        av = pg.find(qn("a:avLst"))
        if av is None:
            av = etree.SubElement(pg, qn("a:avLst"))
        else:
            av.clear()
        gd = etree.SubElement(av, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {round_pct}")
    return shape


def add_cylinder(slide, x, y, w, h, theme_key, border_pt=1.4):
    fill_rgb, border_rgb = THEME[theme_key]
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(border_pt)
    sp = shape.element
    spPr = sp.find(qn("p:spPr"))
    pg = spPr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "can")
        av = pg.find(qn("a:avLst"))
        if av is not None:
            av.clear()
    return shape


def set_text(shape, lines, v_anchor="middle"):
    """
    lines: list of (text, font_size_pt, bold, italic, RGBColor)
    """
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left  = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top   = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    # vertical anchor
    anchor_map = {"top": "t", "middle": "ctr", "bottom": "b"}
    tf._txBody.set("anchor", anchor_map.get(v_anchor, "ctr"))
    for i, (txt, sz, bld, itl, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        r = p.add_run()
        r.text = txt
        r.font.size  = Pt(sz)
        r.font.bold  = bld
        r.font.italic = itl
        r.font.color.rgb = clr


def add_textbox(slide, text, x, y, w, h,
                font_size=8.5, bold=False, italic=False,
                color=C_TITLETXT, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf._txBody.set("anchor", "ctr")
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size  = Pt(font_size)
    r.font.bold  = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_idx(slide, text, x, y):
    """Small index label (e.g. '① Source') above a node."""
    add_textbox(slide, text, x, y, Inches(1.8), Inches(0.18),
                font_size=7.5, color=C_IDXTXT, align=PP_ALIGN.LEFT)


def add_arrow(slide, x1, y1, x2, y2, color=C_LINE, width_pt=1.2, dashed=False):
    """Straight directional connector with an arrowhead at (x2, y2)."""
    lx = min(x1, x2)
    ly = min(y1, y2)
    bw = abs(x2 - x1) or int(Inches(0.01))
    bh = abs(y2 - y1) or int(Inches(0.01))

    cxnSp = etree.SubElement(slide.shapes._spTree, qn("p:cxnSp"))
    nvPr  = etree.SubElement(cxnSp, qn("p:nvCxnSpPr"))
    cNvPr = etree.SubElement(nvPr, qn("p:cNvPr"))
    cNvPr.set("id", str(10000 + len(slide.shapes._spTree)))
    cNvPr.set("name", "arrow")
    etree.SubElement(nvPr, qn("p:cNvCxnSpPr"))
    etree.SubElement(nvPr, qn("p:nvPr"))

    spPr = etree.SubElement(cxnSp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    if x1 > x2:
        xfrm.set("flipH", "1")
    if y1 > y2:
        xfrm.set("flipV", "1")
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(lx)); off.set("y", str(ly))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(bw)); ext.set("cy", str(bh))

    pg = etree.SubElement(spPr, qn("a:prstGeom"))
    pg.set("prst", "line")
    etree.SubElement(pg, qn("a:avLst"))
    etree.SubElement(spPr, qn("a:noFill"))

    ln = etree.SubElement(spPr, qn("a:ln"))
    ln.set("w", str(int(width_pt * 12700)))
    sf = etree.SubElement(ln, qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr"))
    sc.set("val", _rgb_hex(color))

    if dashed:
        pd = etree.SubElement(ln, qn("a:prstDash"))
        pd.set("val", "sysDot")

    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "arrow")
    tail.set("w", "med")
    tail.set("len", "med")


# ─────────────────────────────────────────────────────────────────────────────
# Geometry shortcuts
# ─────────────────────────────────────────────────────────────────────────────
def cx(x, w): return int(x + w / 2)
def cy(y, h): return int(y + h / 2)
def rx(x, w): return int(x + w)       # right edge
def bx(y, h): return int(y + h)       # bottom edge


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

    # background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG

    # ── Title ─────────────────────────────────────────────────────────────
    add_textbox(slide,
                "Process Mining & Agentic Framework — High-Level Pipeline",
                Inches(0.30), Inches(0.06), Inches(12.73), Inches(0.38),
                font_size=15.5, bold=True, color=C_TITLETXT)
    add_textbox(slide,
                ("Source: commitizen-tools/commitizen  ·  "
                 "Object-Centric Process Mining + Declarative Mining + "
                 "LLM-grounded LangGraph Agents"),
                Inches(0.30), Inches(0.45), Inches(12.73), Inches(0.22),
                font_size=9, italic=True, color=C_SUBTXT)

    # ══════════════════════════════════════════════════════════════════════
    # Node dimensions
    # ══════════════════════════════════════════════════════════════════════
    CW, CH = Inches(1.42), Inches(0.66)   # cylinder (DB)
    NW, NH = Inches(1.62), Inches(0.60)   # standard process rect
    MW, MH = Inches(1.90), Inches(0.72)   # mining rect
    LW, LH = Inches(1.90), Inches(0.72)   # LLM rect

    # ══════════════════════════════════════════════════════════════════════
    # Node positions  (all in Inches → converted inline)
    # ══════════════════════════════════════════════════════════════════════

    # ── ① Source ─────────────────────────────────────────────────────────
    N1x, N1y = Inches(0.20), Inches(0.85)
    # back cylinder (stack visual)
    add_cylinder(slide, N1x + Inches(0.06), N1y + Inches(0.06),
                 CW, CH, "source", border_pt=0.8)
    cyl1 = add_cylinder(slide, N1x, N1y, CW, CH, "source")
    set_text(cyl1, [
        ("Event Log",                     9.5, True,  False, C_TITLETXT),
        ("commitizen.json",               7.5, False, True,  C_SUBTXT),
        ("commitizen-tools/commitizen",   6.5, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "① Source", N1x, N1y - Inches(0.21))

    # ── ② Role Determination ─────────────────────────────────────────────
    N2x, N2y = Inches(2.25), Inches(0.87)
    rect2 = add_rect(slide, N2x, N2y, NW, NH, "process")
    set_text(rect2, [
        ("Role",           9.5, True, False, C_TITLETXT),
        ("Determination",  9.5, True, False, C_TITLETXT),
    ])
    add_idx(slide, "② Process", N2x, N2y - Inches(0.21))

    # ── ③ Roles ──────────────────────────────────────────────────────────
    N3x, N3y = Inches(4.55), Inches(0.68)
    N3w, N3h = Inches(1.72), Inches(1.08)
    rect3 = add_rect(slide, N3x, N3y, N3w, N3h, "output")
    set_text(rect3, [
        ("Roles",               9.5, True,  False, C_TITLETXT),
        ("issue_reporter",      7.5, False, True,  C_SUBTXT),
        ("bot",                 7.5, False, True,  C_SUBTXT),
        ("contributor",         7.5, False, True,  C_SUBTXT),
        ("feature_developer",   7.5, False, True,  C_SUBTXT),
        ("maintainer",          7.5, False, True,  C_SUBTXT),
        ("devops_engineer",     7.5, False, True,  C_SUBTXT),
        ("quality_engineer",    7.5, False, True,  C_SUBTXT),
        ("technical_writer",    7.5, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "③ Output", N3x, N3y - Inches(0.21))

    # ── ④ Log Split ───────────────────────────────────────────────────────
    N4x, N4y = Inches(4.55), Inches(2.15)
    rect4 = add_rect(slide, N4x, N4y, NW, NH, "process")
    set_text(rect4, [
        ("Log Split",  9.5, True,  False, C_TITLETXT),
        ("by Role",    8.0, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "④ Process", N4x, N4y - Inches(0.21))

    # ── ⑤ Log Partitions (stacked cylinders) ─────────────────────────────
    N5x, N5y = Inches(6.95), Inches(2.00)
    add_cylinder(slide, N5x + Inches(0.07), N5y + Inches(0.07),
                 CW, CH, "output", border_pt=0.8)
    cyl5 = add_cylinder(slide, N5x, N5y, CW, CH, "output")
    set_text(cyl5, [
        ("Log Partitions",  8.5, True,  False, C_TITLETXT),
        ("Role₁ … Roleₙ",  7.5, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "⑤ Output", N5x, N5y - Inches(0.21))

    # ── ⑥ OCPM / BPMN ────────────────────────────────────────────────────
    N6x, N6y = Inches(9.20), Inches(1.38)
    rect6 = add_rect(slide, N6x, N6y, MW, MH, "mining")
    set_text(rect6, [
        ("OCPM / BPMN",          9.5, True,  False, C_TITLETXT),
        ("Object-Centric &",     7.5, False, True,  C_SUBTXT),
        ("Process Models",       7.5, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "⑥ Mining", N6x, N6y - Inches(0.21))

    # ── ⑦ Declarative Process Mining ─────────────────────────────────────
    N7x, N7y = Inches(9.20), Inches(2.52)
    rect7 = add_rect(slide, N7x, N7y, MW, MH, "mining")
    set_text(rect7, [
        ("Declarative",          9.5, True,  False, C_TITLETXT),
        ("Process Mining",       7.5, False, True,  C_SUBTXT),
        ("DECLARE constraints",  7.0, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "⑦ Mining", N7x, N7y - Inches(0.21))

    # ── ⑧ LLM Process Description ─────────────────────────────────────────
    N8x, N8y = Inches(11.35), Inches(1.38)
    rect8 = add_rect(slide, N8x, N8y, LW, LH, "llm")
    set_text(rect8, [
        ("LLM Process",    9.5, True,  False, C_TITLETXT),
        ("Description",    9.5, True,  False, C_TITLETXT),
        ("(Markdown, NL)", 7.0, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "⑧ Generation", N8x, N8y - Inches(0.21))

    # ── ⑨ Agentic Framework (container) ──────────────────────────────────
    N9x, N9y = Inches(0.20), Inches(3.80)
    N9w, N9h = Inches(12.90), Inches(2.05)
    container = add_rect(slide, N9x, N9y, N9w, N9h, "agent", border_pt=1.8)

    # container header text
    add_textbox(slide,
                "⑨  Agentic Framework — LangGraph Application",
                N9x + Inches(0.12), N9y + Inches(0.06),
                N9w - Inches(0.24), Inches(0.30),
                font_size=11, bold=True,
                color=THEME["agent"][1])   # teal border colour
    add_idx(slide, "⑨ Specification", N9x, N9y - Inches(0.21))

    # sub-agent boxes  — name only, no description
    agent_names = [
        "Issue Reporter\nAgent",
        "Bot Workflow\nAgent",
        "Implementation\nAgent",
        "Quality\nAgent",
        "Technical Writer\nAgent",
    ]
    n_agents = len(agent_names)
    AW = Inches(2.28)
    AH = Inches(1.38)
    total_agents_w = n_agents * AW
    inter_gap = (N9w - Inches(0.30) - total_agents_w) / (n_agents - 1)
    ax0 = N9x + Inches(0.15)
    ay0 = N9y + Inches(0.45)

    for i, name in enumerate(agent_names):
        ax = ax0 + i * (AW + inter_gap)
        arect = add_rect(slide, ax, ay0, AW, AH, "subagent", border_pt=1.1)
        lines = [(part, 10.5, True, False, C_TITLETXT)
                 for part in name.split("\n")]
        set_text(arect, lines)

    # note at bottom of container
    add_textbox(slide,
                ("Grounded in mined role prompts  ·  nodes.py & prompts.py  ·  "
                 "Shared state: IssueOpsState  ·  Policy guard (DECLARE)  ·  "
                 "Human approval gate  ·  GitHub REST (dry-run default)"),
                N9x + Inches(0.15), N9y + N9h - Inches(0.28),
                N9w - Inches(0.30), Inches(0.24),
                font_size=7.5, italic=True,
                color=RGBColor(0x2A, 0x6A, 0x6A))

    # ── ⑩ Evaluation ──────────────────────────────────────────────────────
    N10w, N10h = Inches(2.40), Inches(0.58)
    N10x = int(cx(N9x, N9w) - N10w / 2)
    N10y = Inches(6.25)
    rect10 = add_rect(slide, N10x, N10y, N10w, N10h, "eval")
    set_text(rect10, [
        ("Evaluation",                          10,  True,  False, C_TITLETXT),
        ("qualitative assessment of outputs",   7.5, False, True,  C_SUBTXT),
    ])
    add_idx(slide, "⑩ Evaluation", N10x, N10y - Inches(0.21))

    # ══════════════════════════════════════════════════════════════════════
    # Arrows
    # ══════════════════════════════════════════════════════════════════════

    # ① → ②  (Source → Role Determination, horizontal)
    add_arrow(slide,
              rx(N1x, CW), cy(N1y, CH),
              N2x,         cy(N2y, NH))

    # ① → ④  (Source → Log Split, diagonal)
    add_arrow(slide,
              cx(N1x, CW), bx(N1y, CH),
              cx(N4x, NW), N4y)

    # ② → ③  (Role Determination → Roles)
    add_arrow(slide,
              rx(N2x, NW), cy(N2y, NH),
              N3x,         cy(N3y, N3h))

    # ③ → ④  (Roles → Log Split, vertical — same x centre)
    add_arrow(slide,
              cx(N3x, N3w), bx(N3y, N3h),
              cx(N4x, NW),  N4y)

    # ④ → ⑤  (Log Split → Log Partitions, horizontal)
    add_arrow(slide,
              rx(N4x, NW), cy(N4y, NH),
              N5x,         cy(N5y, CH))

    # ⑤ → ⑥  (Log Partitions → OCPM/BPMN)
    add_arrow(slide,
              rx(N5x, CW), cy(N5y, CH) - Inches(0.12),
              N6x,         cy(N6y, MH))

    # ⑤ → ⑦  (Log Partitions → Declarative Mining)
    add_arrow(slide,
              rx(N5x, CW), cy(N5y, CH) + Inches(0.12),
              N7x,         cy(N7y, MH))

    # ⑥ → ⑧  (OCPM/BPMN → LLM Desc, horizontal)
    add_arrow(slide,
              rx(N6x, MW), cy(N6y, MH),
              N8x,         cy(N8y, LH))

    # ⑥ → ⑨  (OCPM/BPMN → Agentic Framework)
    add_arrow(slide,
              cx(N6x, MW) - Inches(0.20), bx(N6y, MH),
              cx(N6x, MW) - Inches(0.20), N9y)

    # ⑦ → ⑨  (Declarative Mining → Agentic Framework)
    add_arrow(slide,
              cx(N7x, MW), bx(N7y, MH),
              cx(N7x, MW), N9y)

    # ⑧ → ⑨  (LLM Desc → Agentic Framework)
    add_arrow(slide,
              cx(N8x, LW), bx(N8y, LH),
              cx(N8x, LW), N9y)

    # ⑨ → ⑩  (Agentic Framework → Evaluation)
    add_arrow(slide,
              cx(N9x, N9w), bx(N9y, N9h),
              cx(N10x, N10w), N10y)

    # ── Save ──────────────────────────────────────────────────────────────
    out = "langraph-agentic-app/docs/pipeline-diagram.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
